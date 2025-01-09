import streamlit as st
from streamlit_option_menu import option_menu
import base64
from email import header
from html.entities import html5
from importlib.resources import read_binary
import hydralit as hy
from markdown import markdown
from numpy.core.fromnumeric import var
import sys
from streamlit.web import cli as stcli
from PIL import Image
from functions import *
import streamlit.components.v1 as components
import pandas as pd
from st_clickable_images import clickable_images
import numpy as np
import statsmodels.api as sm
from numpy import mean
from numpy import std
from sklearn.datasets import make_classification
from sklearn.model_selection import cross_val_score
from sklearn.model_selection import RepeatedStratifiedKFold
from sklearn.linear_model import LogisticRegression
from matplotlib import pyplot
from sklearn.model_selection import train_test_split
import matplotlib.pyplot as plt
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from sklearn.metrics import silhouette_score
from scipy.spatial.distance import cdist
import seaborn as sns
from io import BytesIO
from statsmodels.formula.api import ols
# from streamlit.state.session_state import SessionState
import tkinter
import matplotlib
# matplotlib.use('TkAgg') testing
# matplotlib.use('Agg')
from sklearn.neural_network import MLPClassifier
from sklearn.metrics import accuracy_score
from mpl_toolkits.mplot3d import Axes3D
from matplotlib import cm
from matplotlib.ticker import LinearLocator, FormatStrFormatter
from sklearn.tree import DecisionTreeRegressor, plot_tree
import sklearn
from sklearn.datasets import make_classification
from sklearn.ensemble import RandomForestClassifier
from sklearn.datasets import make_regression
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import RepeatedKFold
import time
from scipy.stats import pearsonr
from scipy.stats import spearmanr
import dtale
from dtale.views import startup
from dtale.app import get_instance
import webbrowser
import dtale.global_state as global_state
import dtale.app as dtale_app
from matplotlib.pyplot import axis, hist
from scipy import stats as stats
from bioinfokit.analys import stat
from statsmodels.stats.anova import AnovaRM
import statsmodels.api as sm
from statsmodels.graphics.factorplots import interaction_plot
from sklearn.decomposition import PCA
from st_click_detector import click_detector
from streamlit.components.v1 import html
from click_image_dashboards import st_click_image_dashboards
from click_image_advanalytics import st_click_image_advanalytics
from dtale.views import startup
from streamlit_quill import st_quill
import pandas_profiling
from pandas_profiling import ProfileReport
from ydata_profiling import ProfileReport
from streamlit_pandas_profiling import st_profile_report
import dataingestion
from streamlit_card import card
import base64
from io import BytesIO
from sklearn.cluster import KMeans
from sklearn.impute import SimpleImputer
from streamlit_extras.metric_cards import style_metric_cards
from sklearn.preprocessing import LabelEncoder
from tabulate import tabulate
from scipy.stats import mode
from tigerpreds import predictions
import streamlit_antd_components as sac
import pygwalker as pyg
from st_tabs import TabBar
from streamlit_extras.add_vertical_space import add_vertical_space

#add an import to Hydralit
from hydralit import HydraHeadApp
from hydralit import HydraApp

#create a wrapper class
class tigerdashboardshome(HydraHeadApp):

#wrap all your code in this method and you should be done

    def run(self):

        # st.set_page_config(layout="wide")

        menucol1, menucol2, menucol3 = st.columns(3)

        # with menucol1:

            # pagesel = option_menu(None, ["Dashboard", "Advanced analytics", "Make a prediction"],
            #                         icons=['None', 'None', 'None'],
            #                         menu_icon="app-indicator", default_index=0, orientation="horizontal",
            #                         styles={
            #         "container": {"padding": "5!important", "background-color": "#f3e5e5", "color": "#8B0103"},
            #         "icon": {"color": "black", "font-size": "25px"}, 
            #         "nav-link": {"font-size": "12px", "text-align": "left", "margin":"0px", "--hover-color": "#eee"},
            #         "nav-link-selected": {"background-color": "#8B0103", "color": "#f3e5e5", "font-weight": "normal"},
            #     }
            #     )

        if 'edanotes' not in st.session_state:
            st.session_state.edanotes = "Write your notes on EDA here"
        if 'dbnotes' not in st.session_state:
            st.session_state.dbnotes = "Write your notes on dashboard analysis here"

        init = 1
        df_loaded = dataingestion.readdata(init)

        vars = list(df_loaded.columns.values.tolist())

        st.session_state['pagechoice'] = 'dashboards'

        # if pagesel == "Dashboard":

            # value = sac.steps(
            # items=[
            #     sac.StepsItem(title='EDA', description='Explore'),
            #     sac.StepsItem(title='Overall', description='Entire dataset'),
            #     sac.StepsItem(title='Agency/Region', description='Grouped by area'),
            #     sac.StepsItem(title='Manager', description='Manager and below'),
            #     sac.StepsItem(title='Advisor', description='Individual data'),
            # ], format_func='title'
            # )

            # dbrdanalytics, dbrdmetrics, dbrdinteractive, dbrdneeds, dbrdpyg = st.tabs([":red[Analytics]", ":red[Metrics]", ":red[Interactive]", ":red[Needs assessment]", ":red[Pyg]"])

            # dbrdtype = TabBar(tabs=["Analytics","Tab2"],default=0,background = "white", color="red",activeColor="red",fontSize="14px")

            # import extra_streamlit_components as stx

            # chosen_id = stx.tab_bar(data=[
            #     stx.TabBarItemData(id=1, title="ToDo", description="Tasks to take care of"),
            #     stx.TabBarItemData(id=2, title="Done", description="Tasks taken care of"),
            #     stx.TabBarItemData(id=3, title="Overdue", description="Tasks missed out"),
            # ], default=1)

        from st_on_hover_tabs import on_hover_tabs

        st.markdown('<style>' + open('./style.css').read() + '</style>', unsafe_allow_html=True)

        with st.sidebar:
                st.image('indalologo.jpg')
                add_vertical_space(1)
                value = on_hover_tabs(tabName=['Profiling', 'Dashboard/Analytics', 'Validation', 'Data Creation'], 
                                    iconName=['contacts', 'dashboard', 'account_tree', 'pivot_table_chart'],
                                    styles = {'navtab': {'background-color':'#6d0606',
                                                        'color': 'white',
                                                        'font-size': '18px',
                                                        'transition': '.3s',
                                                        'white-space': 'nowrap',
                                                        'text-transform': 'uppercase'},
                                            'tabOptionsStyle': {':hover :hover': {'color': 'red',
                                                                            'cursor': 'pointer'}},
                                            'iconStyle':{'position':'fixed',
                                                            'left':'7.5px',
                                                            'text-align': 'left'},
                                            'tabStyle' : {'list-style-type': 'none',
                                                            'margin-bottom': '30px',
                                                            'padding-left': '30px'}},
                                    key="hoversidebar",
                                    default_choice=0)

        css = '''
        <style>
            .stTabs [data-baseweb="tab-highlight"] {
                background-color:blue;
            }
        </style>
        '''

        st.markdown(css, unsafe_allow_html=True)

        df=df_loaded

        if value=='EDA':
            edalevel = 1
            pass
        if value=='Overall':
            edalevel = 2
            pass
        if value=="Agency/Region":
            edalevel = 3
            pass
        if value=="Manager":
            edalevel = 4
            pass
        if value=='Advisor':
            edalevel = 5
            pass

        if value=='Profiling':

            st.subheader("Dataset profiling")

            with st.spinner("Analyzing and summarizing dataset and generating dataset profile"):

                my_bar = st.progress(0)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                edaenv = st.expander("Guidance on profiling", expanded=False)

                with edaenv:

                    st.info("User guide")
                    @st.cache_resource
                    def show_pdf(file_path):
                        # Opening tutorial from file path
                        with open(file_path, "rb") as f:
                            base64_pdf = base64.b64encode(f.read()).decode('utf-8')

                        # Embedding PDF in HTML
                        pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="1500" height="1000" type="application/pdf"></iframe>'

                        # Displaying File
                        st.markdown(pdf_display, unsafe_allow_html=True)

                    col1, col2,col3= st.columns(3)
                    with col1:  
                        if st.button('Read PDF tutorial',key='1'):            
                            show_pdf('.\Automated flow\DtaleInstructions-compressed.pdf')
                    with col2:
                        st.button('Close PDF tutorial',key='2')                   
                    with col3:
                        with open(".\Automated flow\DtaleInstructions-compressed.pdf", "rb") as pdf_file:
                            PDFbyte = pdf_file.read()
                        st.download_button(label="Download PDF tutorial", key='3',
                                data=PDFbyte,
                                file_name="EDA Instructions.pdf",
                                mime='application/octet-stream')

                datadescrip = st.expander("Description of data")

                with datadescrip:

                    st.write(df.describe(include='all'))
                    
                edaprofiling = st.expander("Profile of dataset", expanded=False)
                
                with edaprofiling:
                
                    profilereport = st.button("Generate profile report", key='profilebutton')
                    if profilereport:
                        description = r"Disclaimer: this profiling report was generated using a sample of the original dataset, conforming to the Central Limit Theorem. Click <a href=\"https://sphweb.bumc.bu.edu/otlt/mph-modules/bs/bs704_probability/BS704_Probability12.html\">here</a> for more information."
                        dataprofile = df.sample(1000).profile_report(dataset={"description": description})
                        dataprofile.config.interactions.targets = ["TenantValue_gender", "TenantValue_maritalstatus", "age", "sum_asset_values", "debtamount", "amount_x"]
                        dataprofile.to_file("TigerDataProfile_" + edalevel + ".pdf")
                        st.success("Dataset profile saved to file 'TigerDataProfile_" + edalevel + ".pdf")
                    # build_profile()

            dtaleanalysis = st.button("Perform detailed EDA", key='dtalebutton')

            if dtaleanalysis:

                startup(data_id="1", data=df) # All records, no OHE

                if get_instance("1") is None:
                    startup(data_id="1", data=df)

                d=get_instance("1")

                html = f"""<iframe src="/dtale/main/1" height="1000" width="1400"></iframe>""" # Iframe
    
                st.markdown(html, unsafe_allow_html=True)

                st.session_state.corr_img = d.get_corr_matrix()
                st.session_state.corr_df = d.get_corr_matrix(as_df=True)
                st.session_state.pps_img = d.get_pps_matrix()
                st.session_state.pps_df = d.get_pps_matrix(as_df=True)

                print(st.session_state.corr_df)

            # Spawn a new Quill editor
            st.subheader("Notes on EDA analysis")
            edacontent = st_quill(placeholder="Write your notes here", value=st.session_state.edanotes, key="edaquill")

            st.session_state.edanotes = edacontent

            st.write("Exploratory data analysis took ", time.time() - start_time, "seconds to run")

        if value == 'Dashboard/Analytics':

            dbrdanalytics, dbrdmetrics, dbrdinteractive, dbrdneeds, dbrdpyg = st.tabs([":red[Analytics]", ":red[Metrics]", ":red[Interactive]", ":red[Needs assessment]", ":red[Custom visualization]"])

            st.subheader("Tiger analytics")

            with st.spinner("Analyzing dataset and generating dataset statistics"):

                my_bar = st.progress(0)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

            # Helper function for section headers
            def section_header(header_text, header_size=2):
                header_tag = f"h{header_size}"
                centered_style = 'text-align: center;'
                st.markdown(f"<{header_tag} style='{centered_style}'>{header_text}</{header_tag}>", unsafe_allow_html=True)

            # Apply CSS styling to create a rounded block format #3182A8
            image_block_style = """
                display: flex;
                justify-content: center;
                align-items: center;
                padding: 20px;
                border-radius: 10px;
                background-color: #E8E8E8;
                box-shadow: 0 2px 5px rgba(0, 0, 0, 0.1);
                margin-bottom: 20px; /* Add margin bottom to create space below the block */
            """

            def style_custom_metric_cards(
                background_color: str = "#E8E8E8",
                border_size_px: int = 1,
                border_color: str = "#CCC",
                border_radius_px: int = 5,
                border_left_color: str = "#207DCE",
                box_shadow: bool = True,
            ):
                box_shadow_str = (
                    "box-shadow: 0 0.15rem 1.75rem 0 rgba(58, 59, 69, 0.15) !important;"
                    if box_shadow
                    else "box-shadow: none !important;"
                )
                st.markdown(
                    f"""
                    <style>
                        div[data-testid="custom-metric-container"] {{
                            background-color: {background_color};
                            border: {border_size_px}px solid {border_color};
                            padding: 1.5% 5% 5% 10%;
                            margin-top: -12px;
                            border-radius: {border_radius_px}px;
                            border-left: 0.5rem solid {border_left_color} !important;
                            {box_shadow_str}
                        }}
                        div[data-testid="custom-metric-container"] .metric-value {{
                            font-size: 36px;
                        }}
                        div[data-testid="custom-metric-container"] .metric-label {{
                            font-size: 14px;
                        }}
                        div[data-testid="custom-metric-container"] .metric-delta {{
                            font-size: 14px;
                            text-align: left;
                            white-space: pre-wrap;
                        }}
                    </style>
                    """,
                    unsafe_allow_html=True,
                )

            with dbrdanalytics:

                demogoverview = st.expander("Demographics overview", expanded=True)

                with demogoverview:

                    # Initial dashboard EDA

                    section_header("Demographics overview", header_size=3)

                    def generate_card_with_overlay(image_url, button1_text, button2_text, button3_text, card_text, expln_text, card_title):
                        html_code = f'''
                            <style>
                                .card-container {{
                                    display: flex;
                                    border: 1px solid #fff; /* Change to white background */
                                    width: 375px;
                                    overflow: hidden;
                                    flex-direction: column;
                                    margin-bottom: 20px;
                                }}
                                .card-left {{
                                    padding: 10px;
                                    cursor: pointer;
                                    font-family: 'Roboto', sans-serif;
                                    font-size: 12px;
                                    display: flex;
                                    flex-direction: column;
                                    align-items: center;
                                    justify-content: flex-start;
                                }}
                                .card-left img {{
                                    max-width: 100%;
                                    max-height: 100%;
                                }}
                                .card-right {{
                                    padding: 10px;
                                    font-family: 'Roboto', sans-serif;
                                    font-size: 14px;
                                    overflow-y: hidden;
                                }}
                                .overlay {{
                                    display: none;
                                    position: fixed;
                                    top: 1px;
                                    left: 0;
                                    width: 90%;
                                    height: 90%;
                                    background-color: transparent; /* Transparent background */
                                    z-index: 1;
                                    justify-content: center;
                                    align-items: center;
                                }}
                                .overlay-image {{
                                    max-width: 80%; /* Adjust the width of the overlay image */
                                    max-height: 80%; /* Adjust the height of the overlay image */
                                }}
                                .card-title {{
                                    text-align: center;
                                    margin-top: 10px;
                                    font-family: 'Roboto', sans-serif;
                                    font-size: 14px;
                                }}
                                .button-container {{
                                    display: flex;
                                    justify-content: center;
                                }}
                                button {{
                                    background-color: #207DCE;
                                    color: white;
                                    padding: 10px 15px;
                                    border: none;
                                    border-radius: 15px; /* Make buttons rounded */
                                    cursor: pointer;
                                    margin: 5px;
                                }}
                            </style>
                            <div class="card-container" style="font-family: 'Roboto', sans-serif;">
                                <h3 class="card-title">{card_title}</h3>
                                <div class="card-left" onclick="openOverlay()">
                                    <img src="{image_url}" style="width: 375px; height: 250px; margin-top: 25px;">
                                </div>
                                <div class="card-right">
                                    <div class="button-container">
                                        <button id="graph_expln_button">{button1_text}</button>
                                        <button id="graph_expln_button2">{button2_text}</button>
                                        <button id="reset_button">{button3_text}</button>
                                    </div>
                                    <p id="dynamic_heading"></p>
                                    <p class="card-text" id="dynamic_text">{" "}</p>
                                    <p class="card-text"><small class="text-muted">Last updated 3 mins ago</small></p>
                                </div>
                            </div>
                            <div class="overlay" id="imageOverlay" onclick="closeOverlay()">
                                <img src="{image_url}" class="overlay-image">
                            </div>
                            <script>
                                function openOverlay() {{
                                    document.getElementById("imageOverlay").style.display = "block";
                                }}
                                function closeOverlay() {{
                                    document.getElementById("imageOverlay").style.display = "none";
                                }}
                                document.getElementById("graph_expln_button").addEventListener("click", function() {{
                                    document.getElementById("dynamic_heading").style.fontWeight = "bold";
                                    document.getElementById("dynamic_heading").innerText = "Analysis";
                                    document.getElementById("dynamic_text").innerText = "{card_text}";
                                    adjustCardHeight();
                                }});
                                document.getElementById("graph_expln_button2").addEventListener("click", function() {{
                                    document.getElementById("dynamic_heading").style.fontWeight = "bold";
                                    document.getElementById("dynamic_heading").innerText = "Recommendations";
                                    document.getElementById("dynamic_text").innerText = "{expln_text}";
                                    adjustCardHeight();
                                }});
                                document.getElementById("reset_button").addEventListener("click", function() {{
                                    document.getElementById("dynamic_heading").innerText = "";
                                    document.getElementById("dynamic_text").innerText = "";
                                    adjustCardHeight();
                                }});
                                function adjustCardHeight() {{
                                    var cardContainer = document.querySelector('.card-container');
                                    cardContainer.style.height = 'auto';
                                }}
                            </script>
                        '''
                        return html_code
                    
                    colm1, colm2, colm3 = st.columns(3)

                    with colm1:
                        image_url = 'suppliertype.png'
                        buffer = BytesIO()
                        with open(image_url, 'rb') as image_file:
                        # Read the file's content
                            image_data = image_file.read()
                        encoded_data = base64.b64encode(image_data).decode('utf-8')
                        image_url1 = f'data:image/png;base64,{encoded_data}'
                        card_title1 = "Count of Supplier by type"
                        button1_text1 = "Analysis"
                        button2_text1 = "Recommendations"
                        button3_text1 = "Clear"
                        card_text1 = "There are 20 possible supplier types, with the dominant being 'Supplier'. This might be a misclassfication."
                        expln_text1 = "We recommend you recode supplier types to alleviate non-specific categorization."
                        html_code1 = generate_card_with_overlay(image_url1, button1_text1, button2_text1, button3_text1, card_text1, expln_text1, card_title1)
                        from streamlit.components.v1 import html
                        html(f'<div id="card1_Wrapper">{html_code1}</div>', width=1200, height=550)  # Adjust the overall width and height of the HTML element

                    # Add JavaScript to make the row height dynamic
                    st.markdown(
                        """
                        <script>
                            const card1Wrapper = document.getElementById('card1_wrapper');

                            function setDynamicHeight() {
                                const cardHeight = card1Wrapper.scrollHeight;
                                card1Wrapper.style.height = cardHeight + 'px';
                            }

                            // Call the function when the page loads
                            setDynamicHeight();

                            // Call the function whenever the content inside the card changes
                            card1Wrapper.addEventListener('DOMSubtreeModified', setDynamicHeight);
                        </script>
                        """,
                        unsafe_allow_html=True
                    )
            
            with dbrdmetrics:

                from streamlit_extras.metric_cards import style_metric_cards

                headlinemetrics = st.expander("Headline metrics", expanded=True)

                # Apply custom styles to custom metric cards
                style_custom_metric_cards()

                with headlinemetrics:
                    
                    # # Section 1: Count of Records

                    section_header("Count of Usable Records", header_size=3)
                    style_metric_cards(background_color='#E8E8E8', border_left_color='#207DCE', box_shadow=True)
                    col1, col2, col3 = st.columns(3)
                    record_count = len(df)
                    # formatted_count = f"{record_count:,.0f}"
                    formatted_count = f"{record_count:,.0f}"
                    formatted_count = formatted_count.replace(',', ' ')
                    col2.metric(label="Count of Usable Records", value=formatted_count)

                    # Section 2: Demographic Distribution

                    section_header("Demographic Distribution", header_size=3)
                    col1, col2, col3 = st.columns(3)
                    gender_distribution_male = 1
                    gender_distribution_female = 2
                    total = gender_distribution_male + gender_distribution_female
                    marital_status_distribution = 6
                    # Define the values for the custom card
                    with col1:
                        label = 'Gender distribution %'
                        value = "Males: "
                        delta = "Females: "
                        # Display the custom metric card with custom content
                        st.markdown(
                            f"""
                            <div data-testid="custom-metric-container">
                                <div class="metric-label">{label}</div>
                                <div class="metric-value">{value}</div>
                                <div class="custom-metric-delta">{delta}</div>
                            </div>
                            """,
                            unsafe_allow_html=True
                        )

            with dbrdpyg:

                from pygwalker.api.streamlit import StreamlitRenderer, init_streamlit_comm

                tableauboard = st.expander('Custom visualizations', expanded=True)

                with tableauboard:

                    df_tableau= df.copy()

                    # df_tableau = np.ascontiguousarray(df_tableau)

                    print(df_tableau.head(5))

                    for col in df_tableau.select_dtypes(include=['category']).columns:
                        if "unknown" not in df_tableau[col].cat.categories:
                            df_tableau[col].cat.add_categories("unknown", inplace=True)

                    # For numeric columns
                    for col in df_tableau.select_dtypes(include=['float64', 'int64']).columns:
                        df_tableau[col] = df_tableau[col].fillna(0)

                    # For non-numeric columns
                    for col in df_tableau.select_dtypes(exclude=['float64', 'int64']).columns:
                        df_tableau[col] = df_tableau[col].fillna("unknown")
                    
                    from pygwalker.api.streamlit import StreamlitRenderer, init_streamlit_comm

                    # Establish communication between pygwalker and streamlit
                    init_streamlit_comm()
                    
                    # Get an instance of pygwalker's renderer. You should cache this instance to effectively prevent the growth of in-process memory.
                    @st.cache_resource
                    def get_pyg_renderer() -> "StreamlitRenderer":
                        # When you need to publish your app to the public, you should set the debug parameter to False to prevent other users from writing to your chart configuration file.
                        return StreamlitRenderer(df_tableau, spec="./gw_config.json", debug=False)
                    
                    renderer = get_pyg_renderer()
                    
                    # # Render your data exploration interface. Developers can use it to build charts by drag and drop.
                    renderer.render_explore()

                    pyg_app = StreamlitRenderer(df_tableau)

                    # pyg_app.explorer()

                    # pass

            # Spawn a new Quill editor
            st.subheader("Notes on dashboard data analysis")
            dbcontent = st_quill(placeholder="Write your notes here", value=st.session_state.dbnotes, key="dbquill")

            st.session_state.dbanotes = dbcontent

            st.write("Exploratory data analysis took ", time.time() - start_time, "seconds to run")

        if value == 'Validation':

            st.subheader("Data validation")

            with st.spinner("Setting up data validation"):

                my_bar = st.progress(0)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

            from mitosheet.streamlit.v1 import spreadsheet
            spreadsheet(df)

        if value == 'Data Creation':

            st.subheader("Data creation")

            with st.spinner("Data creation form opening"):

                my_bar = st.progress(0)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                st.success("Data validation rules are baked into the data creation form")

            webbrowser.open_new_tab('https://ee.kobotoolbox.org/x/NwleEkrR')

        else:

            st.warning('Please select a dashboard option above')

        with st.expander("Build results presentation", expanded=False):

            st.info("When finished, click below to build a presentation with your results")

            buildpres = st.button("Build presentation")

            if buildpres == True:

                # Building Powerpoint presentation

                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.dml.color import RGBColor
                from pptx.util import Inches, Pt
                from pptx.enum.dml import MSO_THEME_COLOR
                title='   Analytics Playground\n\
                Results from analysis'
                APlogo='./Powerpoint/APlogo.png'
                ABIlogo='./Powerpoint/ABIlogo.png'
                prs = Presentation()

                # Slide 1

                # Add colour bar

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, Inches(9/1.5),Inches(16),Inches(9/8.5)
                )
                shape.shadow.inherit = False
                fill=shape.fill
                fill.solid()
                fill.fore_color.rgb=RGBColor(255,0,0)
                shape.text= title
                line=shape.line
                line.color.rgb=RGBColor(255,0,0)
                logo1=slide.shapes.add_picture(APlogo,Inches(13.5),Inches(6.0),height=Inches(1.08),width=Inches(1.0))
                logo2=slide.shapes.add_picture(ABIlogo,Inches(14.5),Inches(5.8),height=Inches(1.5),width=Inches(1.51))

                # Add text box for results

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.3))
                shape.shadow.inherit = False
                fill=shape.fill
                fill.solid()
                fill.fore_color.rgb=RGBColor(255,0,0)
                shape.text= "   Results from Exploratory Data Analysis"
                line=shape.line
                line.color.rgb=RGBColor(255,0,0)
                logo1=slide.shapes.add_picture(APlogo,Inches(14.5),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
                logo2=slide.shapes.add_picture(ABIlogo,Inches(15.0),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
                left = Inches(1)
                top = Inches(2)
                width = Inches(5)
                height = Inches(5)
                text_box=slide.shapes.add_textbox(left, top, width, height)
                tb=text_box.text_frame
                tb.text = st.session_state.edanotes
                prg=tb.add_paragraph()
                prg.text=" "
                prg=tb.add_paragraph()
                prg.text=''
                correlpic = slide.shapes.add_picture('correl.jpg', Inches(8), Inches(1.3), height=Inches(3.7), width=Inches(6.3))
                ppspic = slide.shapes.add_picture('pps.jpg', Inches(8), Inches(5.1), height=Inches(3.7), width=Inches(7.3))

                prs.save('EDA_presentation.pptx')

                os.startfile("EDA_presentation.pptx")

        st.markdown("""<link href="https://cdn.jsdelivr.net/npm/bootstrap@5.0.2/dist/css/bootstrap.min.css" rel="stylesheet" integrity="sha384-EVSTQN3/azprG1Anm3QDgpJLIm9Nao0Yz1ztcQTwFspd3yD65VohhpuuCOmLASjC" crossorigin="anonymous">""", unsafe_allow_html=True)
        st.markdown("""<script src="https://cdn.jsdelivr.net/npm/bootstrap@5.0.2/dist/js/bootstrap.bundle.min.js" integrity="sha384-MrcW6ZMFYlzcLA8Nl+NtUVF0sA7MsXsP1UyJoMp4YLEuNSfAP+JcXn/tWtIaxVXM" crossorigin="anonymous"></script>""", unsafe_allow_html=True)

        image_block_style2 = """
            display: flex;
            justify-content: center;
            align-items: center;
            padding: 10px;
            border-radius: 10px;
            background-color: #E8E8E8;
            box-shadow: 0 2px 5px rgba(0, 0, 0, 0.1);
            margin-bottom: 10px; /* Add margin bottom to create space below the block */
        """

        # if pagesel == 'Advanced analytics':

        #     df = df_loaded

        #     import os
        #     import openai
        #     from dotenv import load_dotenv
        #     from pandasai import SmartDataframe
        #     from pandasai.llm.openai import OpenAI

        #     load_dotenv()
        #     openai.api_key = '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' # os.environ["OPENAI_API_KEY"]
        #     llm = OpenAI(api_token=openai.api_key)
        #     st.title("Data analysis")
        #     sdf = SmartDataframe(df, config={"llm":llm})
        #     prompt = st.text_area("Enter your prompt")
        #     if st.button("Generate"):
        #         if prompt:
        #             with st.spinner("Generating response"):
        #                 response = sdf.chat(prompt)
        #                 st.success(response)
        #                 st.set_option('deprecation.showPyplotGlobalUse', False)
        #                 st.pyplot()
        #         else:
        #             st.warning("Please enter a prompt")

